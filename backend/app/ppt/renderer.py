"""
PPT Renderer — Programmatic PPTX generation from Slide DSL.

Rule: the renderer is NOT an LLM. It follows template constraints.
"""
import asyncio
import io
import logging
import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.core.config import settings

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)
if not logger.handlers:
    _h = logging.StreamHandler()
    _h.setFormatter(logging.Formatter('%(levelname)s [ppt] %(message)s'))
    logger.addHandler(_h)

# ---- Image fetching via Tavily ----
_IMAGE_CACHE: dict[str, bytes] = {}
_MAX_IMAGE_CACHE = 64


async def _fetch_images(queries: list[str]) -> list[bytes]:
    """Search images via Tavily and download them in parallel."""
    if not settings.tavily_api_key:
        return await _fallback_images(queries)

    async def _fetch_one(q: str) -> bytes | None:
        if q in _IMAGE_CACHE:
            logger.debug("Image cache hit: %s", q[:60])
            return _IMAGE_CACHE[q]
        try:
            from tavily import TavilyClient
            client = TavilyClient(api_key=settings.tavily_api_key)
            resp = client.search(query=q, search_depth="basic", include_images=True, max_results=3)
            image_urls = resp.get("images", [])
            if image_urls:
                logger.info("Tavily found %d images for query '%s'", len(image_urls), q[:60])
                img_bytes = await _download_image(image_urls[0])
                if img_bytes:
                    _cache_image(q, img_bytes)
                    logger.info("Downloaded image %d bytes from %s", len(img_bytes), image_urls[0][:80])
                    return img_bytes
                else:
                    logger.warning("Failed to download image from %s", image_urls[0][:80])
            else:
                logger.warning("Tavily returned no images for query '%s'", q[:60])
        except Exception:
            logger.exception("Tavily search failed for query '%s'", q[:60])
        return None

    results = await asyncio.gather(*(_fetch_one(q) for q in queries[:3]))
    return [r for r in results if r is not None]


def _cache_image(key: str, data: bytes) -> None:
    if len(_IMAGE_CACHE) >= _MAX_IMAGE_CACHE:
        _IMAGE_CACHE.pop(next(iter(_IMAGE_CACHE)))
    _IMAGE_CACHE[key] = data


async def _download_image(url: str) -> bytes | None:
    try:
        async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
            resp = await client.get(url)
            if resp.status_code == 200 and len(resp.content) > 2000:
                return resp.content
    except Exception:
        pass
    return None


async def _fallback_images(queries: list[str]) -> list[bytes]:
    logger.info("Using Unsplash fallback for images (no Tavily API key)")
    async def _fetch_one(q: str) -> bytes | None:
        if q in _IMAGE_CACHE:
            return _IMAGE_CACHE[q]
        try:
            url = f"https://source.unsplash.com/800x600/?{q}"
            async with httpx.AsyncClient(timeout=10, follow_redirects=True) as client:
                resp = await client.get(url)
                if resp.status_code == 200 and len(resp.content) > 2000:
                    _cache_image(q, resp.content)
                    logger.info("Unsplash download success for '%s': %d bytes", q[:60], len(resp.content))
                    return resp.content
                else:
                    logger.warning("Unsplash returned status %d, %d bytes for '%s'",
                                   resp.status_code, len(resp.content), q[:60])
        except Exception:
            logger.exception("Unsplash download failed for '%s'", q[:60])
        return None

    results = await asyncio.gather(*(_fetch_one(q) for q in queries[:2]))
    return [r for r in results if r is not None]


def _embed_image(slide, image_bytes: bytes, left, top, width, height):
    try:
        slide.shapes.add_picture(io.BytesIO(image_bytes),
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
        logger.info("Embedded image at (%.1f, %.1f) size %.1fx%.1f, %d bytes",
                     left, top, width, height, len(image_bytes))
    except Exception as e:
        logger.warning("Failed to embed image (%d bytes), falling back to placeholder: %s",
                       len(image_bytes), e)
        rect = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(height))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _color("E8E8E8")
        rect.line.fill.background()


# ---- Theme System ----
THEMES = {
    "japan": {"bg": "F5F0E8", "accent": "9B1B30", "text": "2D2D2D", "sub": "8B7355"},
    "minimal": {"bg": "FFFFFF", "accent": "111111", "text": "1A1A1A", "sub": "757575"},
    "tropical": {"bg": "FFF8F0", "accent": "E8734A", "text": "2C3E50", "sub": "8E6E53"},
    "ocean": {"bg": "F0F4F8", "accent": "1B4F72", "text": "1C2833", "sub": "5D6D7E"},
    "editorial": {"bg": "FAFAFA", "accent": "2C3E50", "text": "1A1A1A", "sub": "7F8C8D"},
    "nature": {"bg": "F4F6F0", "accent": "3E6B48", "text": "2D3436", "sub": "6B7F6B"},
}

DEFAULT_THEME = THEMES["minimal"]

WIDTH = Inches(13.333)
HEIGHT = Inches(7.5)


def _color(hex_str: str) -> RGBColor:
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


async def render_pptx(slides: list[dict], theme_name: str = "minimal",
                     assets: list[dict] | None = None) -> bytes:
    prs = Presentation()
    prs.slide_width = WIDTH
    prs.slide_height = HEIGHT

    theme = THEMES.get(theme_name, DEFAULT_THEME)
    assets = assets or []

    blank_layout = prs.slide_layouts[6]

    # Pre-fetch images for ALL slides that have asset queries
    image_futures = {}
    for idx, s in enumerate(slides):
        queries = _get_asset_queries(assets, idx)
        if queries:
            image_futures[idx] = asyncio.ensure_future(_fetch_images(queries))

    for idx, s in enumerate(slides):
        slide_type = s.get("slide_type", "content")
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, theme["bg"])

        images = None
        if idx in image_futures:
            images = await image_futures[idx]

        render = _SLIDE_RENDERERS.get(slide_type, _render_content)
        render(slide, s, theme, images or [])

    # Log image embedding summary
    slides_with_images = [i for i in image_futures]
    logger.info("PPT render done: %d slides, %d with image queries, assets=%s",
                 len(slides), len(slides_with_images),
                 [(a.get("slide_index"), a.get("queries", [])[:1]) for a in assets])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _get_asset_queries(assets: list[dict], slide_index: int) -> list[str]:
    for a in assets:
        if a.get("slide_index") == slide_index:
            return a.get("queries", [])
    return []


_SLIDE_RENDERERS = {}  # populated after all render functions are defined


def _set_bg(slide, color: str):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _color(color)


def _add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                  color="2D2D2D", align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _color(color)
    p.font.name = font_name
    p.alignment = align
    return tf


# ---- Slide Renderers ----

def _render_cover(slide, s, theme, images: list[bytes] | None = None):
    images = images or []
    title = s.get("title", "Travel Plan")
    subtitle = s.get("subtitle", "")

    if images:
        _embed_image(slide, images[0], 6.8, 0, 6.6, 7.5)
        panel = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(7.2), Inches(7.5))
        panel.fill.solid()
        panel.fill.fore_color.rgb = _color(theme["bg"])
        panel.line.fill.background()

    _add_text_box(slide, 0.8, 2.2, 5.5, 1.5, title, 44, True, theme["accent"],
                  PP_ALIGN.LEFT, "Arial Black")
    if subtitle:
        _add_text_box(slide, 0.8, 3.7, 5, 0.8, subtitle, 20, False, theme["sub"])
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(4.7), Inches(3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = _color(theme["accent"])
    line.line.fill.background()


def _render_timeline(slide, s, theme, images: list[bytes] | None = None):
    title = s.get("title", "Itinerary")
    subtitle = s.get("subtitle", "")
    days = s.get("days", [])

    if images:
        _embed_image(slide, images[0], 9.8, 0.3, 3.2, 2.4)

    _add_text_box(slide, 0.8, 0.4, 8.5, 0.7, title, 28, True, theme["text"])
    if subtitle:
        _add_text_box(slide, 0.8, 1.0, 11, 0.5, subtitle, 16, False, theme["sub"])

    n = len(days)
    if n == 0:
        return
    col_w = 11.5 / n
    for i, day in enumerate(days):
        x = 0.8 + i * col_w
        _add_text_box(slide, x, 1.8, col_w - 0.3, 0.4,
                      f"Day {day.get('day', i+1)}", 18, True, theme["accent"])
        _add_text_box(slide, x, 2.2, col_w - 0.3, 0.8,
                      day.get("theme", ""), 14, False, theme["sub"])
        activities = day.get("activities", [])
        y = 3.1
        for act in activities[:4]:
            _add_text_box(slide, x, y, col_w - 0.3, 0.6,
                          f"• {act[:60]}", 11, False, theme["text"])
            y += 0.7


def _render_cards(slide, s, theme, images: list[bytes] | None = None):
    title = s.get("title", "")
    items = s.get("items", [])

    if images:
        _embed_image(slide, images[0], 9.8, 0.3, 3.2, 2.4)

    _add_text_box(slide, 0.8, 0.4, 8.5, 0.7, title, 28, True, theme["text"])

    normalized = []
    for item in items:
        if isinstance(item, str):
            normalized.append({"name": item[:60], "desc": ""})
        elif isinstance(item, dict):
            normalized.append({
                "name": str(item.get("name", item.get("title", "")))[:60],
                "desc": str(item.get("desc", item.get("subtitle", "")))[:120]
            })
    items = normalized

    n = len(items)
    if n == 0:
        return
    cols = 2 if n <= 4 else 3
    col_w = 11.5 / cols
    for i, item in enumerate(items):
        row = i // cols
        col = i % cols
        x = 0.8 + col * col_w
        y = 1.6 + row * 2.5
        card = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(col_w - 0.4), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = _color("FFFFFF")
        card.line.color.rgb = _color("E0E0E0")
        card.line.width = Pt(0.5)
        _add_text_box(slide, x + 0.2, y + 0.15, col_w - 0.8, 0.4,
                      item["name"][:40], 16, True, theme["accent"])
        _add_text_box(slide, x + 0.2, y + 0.7, col_w - 0.8, 1.0,
                      item["desc"][:120], 11, False, theme["text"])


def _render_budget(slide, s, theme, images: list[bytes] | None = None):
    title = s.get("title", "Budget")
    items = s.get("items", [])
    total = s.get("total", "")
    _add_text_box(slide, 0.8, 0.4, 8.5, 0.7, title, 28, True, theme["text"])

    headers = ["类别", "明细", "金额"]
    col_widths = [2.5, 6.0, 3.0]
    x_positions = [0.8]
    for w in col_widths[:-1]:
        x_positions.append(x_positions[-1] + w)

    y = 1.6
    for j, h in enumerate(headers):
        _add_text_box(slide, x_positions[j], y, col_widths[j], 0.4,
                      h, 13, True, theme["accent"])
    y += 0.5
    for item in items[:8]:
        if isinstance(item, str):
            _add_text_box(slide, 0.8, y, 11, 0.35, f"• {item[:80]}", 12, False, theme["text"])
        else:
            for j, key in enumerate(["category", "detail", "amount"]):
                _add_text_box(slide, x_positions[j], y, col_widths[j], 0.35,
                              str(item.get(key, ""))[:50], 12, False, theme["text"])
        y += 0.4

    if total:
        _add_text_box(slide, 0.8, y + 0.2, 11, 0.4, total, 16, True, theme["accent"])


def _render_map_slide(slide, s, theme, images: list[bytes] | None = None):
    images = images or []
    title = s.get("title", "Route Map")
    _add_text_box(slide, 0.8, 0.4, 11, 0.7, title, 28, True, theme["text"])

    if images:
        _embed_image(slide, images[0], 0.8, 1.5, 11.5, 5.5)
    else:
        _add_text_box(slide, 0.8, 2.5, 11.5, 4.0,
                      s.get("description", "路线示意图"),
                      14, False, theme["sub"])


def _render_ending(slide, s, theme, images: list[bytes] | None = None):
    _add_text_box(slide, 1.5, 2.8, 10, 1.0, s.get("title", "Bon Voyage! ✈️"),
                  42, True, theme["accent"], PP_ALIGN.CENTER, "Arial Black")
    _add_text_box(slide, 1.5, 3.8, 10, 0.6, s.get("subtitle", "Have a wonderful trip"),
                  20, False, theme["sub"], PP_ALIGN.CENTER)


def _render_content(slide, s, theme, images: list[bytes] | None = None):
    title = s.get("title", "")
    bullets = s.get("bullets", [])
    _add_text_box(slide, 0.8, 0.4, 11, 0.7, title, 28, True, theme["text"])
    y = 1.6
    for b in bullets[:8]:
        _add_text_box(slide, 1.2, y, 10.5, 0.5, f"• {b[:100]}", 14, False, theme["text"])
        y += 0.55


_SLIDE_RENDERERS = {
    "cover": _render_cover,
    "map": _render_map_slide,
    "timeline": _render_timeline,
    "overview": _render_timeline,
    "itinerary": _render_timeline,
    "food": _render_cards,
    "hotel": _render_cards,
    "tips": _render_cards,
    "budget": _render_budget,
    "ending": _render_ending,
}
